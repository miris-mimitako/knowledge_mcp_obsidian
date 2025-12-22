"""
ファイルパーサーモジュール
各種ファイル形式（PDF、Word、PowerPoint、Excel、Text/Markdown）の解析
"""
import os
from pathlib import Path
from typing import List, Dict, Optional
import re


class FileParser:
    """各種ファイル形式を解析するベースクラス"""
    
    @staticmethod
    def should_skip_excel_sheet(cells_text: str, min_text_ratio: float = 0.3) -> bool:
        """
        Excelシートが数値のみで自然言語が含まれないかを判定
        
        Args:
            cells_text: セルのテキスト内容を連結した文字列
            min_text_ratio: テキストが占める最小の割合（0.0-1.0）
        
        Returns:
            スキップすべき場合はTrue
        """
        if not cells_text or len(cells_text.strip()) == 0:
            return True
        
        # 数字、記号、空白のみのセルをカウント
        numeric_pattern = re.compile(r'^[\d\s\.\,\-\+\%\(\)\/]+$')
        lines = cells_text.split('\n')
        
        text_lines = 0
        total_lines = 0
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            total_lines += 1
            # 数値のみでない行があればテキスト行としてカウント
            if not numeric_pattern.match(line):
                text_lines += 1
        
        if total_lines == 0:
            return True
        
        text_ratio = text_lines / total_lines
        return text_ratio < min_text_ratio


class PDFParser(FileParser):
    """PDFファイルパーサー"""
    
    @staticmethod
    def parse(file_path: str) -> List[Dict[str, str]]:
        """
        PDFファイルをページ単位で解析
        
        Args:
            file_path: PDFファイルのパス
        
        Returns:
            各ページのテキストと位置情報のリスト
        """
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise ImportError("PyMuPDF (fitz) がインストールされていません。pip install pymupdf を実行してください。")
        
        results = []
        doc = fitz.open(file_path)
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            
            if text.strip():
                results.append({
                    "content": text,
                    "location_info": f"Page {page_num + 1}"
                })
        
        doc.close()
        return results


class WordParser(FileParser):
    """Word (.docx) ファイルパーサー"""
    
    @staticmethod
    def parse(file_path: str) -> List[Dict[str, str]]:
        """
        Wordファイルをページ単位で解析（段落単位で分割）
        
        Args:
            file_path: Wordファイルのパス
        
        Returns:
            各セクションのテキストと位置情報のリスト
        """
        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx がインストールされていません。pip install python-docx を実行してください。")
        
        results = []
        doc = Document(file_path)
        
        current_section = []
        page_num = 1
        
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                current_section.append(text)
        
        # 全テキストを1つのセクションとして保存
        if current_section:
            full_text = '\n'.join(current_section)
            results.append({
                "content": full_text,
                "location_info": f"Page {page_num}"
            })
        
        return results if results else []


class PowerPointParser(FileParser):
    """PowerPoint (.pptx) ファイルパーサー"""
    
    @staticmethod
    def parse(file_path: str) -> List[Dict[str, str]]:
        """
        PowerPointファイルをスライド単位で解析
        
        Args:
            file_path: PowerPointファイルのパス
        
        Returns:
            各スライドのテキストと位置情報のリスト
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx がインストールされていません。pip install python-pptx を実行してください。")
        
        results = []
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_texts.append(shape.text.strip())
            
            if slide_texts:
                full_text = '\n'.join(slide_texts)
                results.append({
                    "content": full_text,
                    "location_info": f"Slide {slide_num}"
                })
        
        return results


class ExcelParser(FileParser):
    """Excel (.xlsx) ファイルパーサー"""
    
    @staticmethod
    def parse(file_path: str) -> List[Dict[str, str]]:
        """
        Excelファイルをシート単位で解析（数値のみのシートは除外）
        
        Args:
            file_path: Excelファイルのパス
        
        Returns:
            各シートのテキストと位置情報のリスト
        """
        try:
            from openpyxl import load_workbook
        except ImportError:
            raise ImportError("openpyxl がインストールされていません。pip install openpyxl を実行してください。")
        
        results = []
        wb = load_workbook(file_path, data_only=True)
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            cells_text = []
            
            for row in sheet.iter_rows(values_only=True):
                row_text = []
                for cell in row:
                    if cell is not None:
                        cell_str = str(cell).strip()
                        if cell_str:
                            row_text.append(cell_str)
                
                if row_text:
                    cells_text.append(' '.join(row_text))
            
            full_text = '\n'.join(cells_text)
            
            # 数値のみのシートは除外
            if not FileParser.should_skip_excel_sheet(full_text):
                results.append({
                    "content": full_text,
                    "location_info": f"Sheet: {sheet_name}"
                })
        
        wb.close()
        return results


class TextParser(FileParser):
    """テキストファイルパーサー（.txt, .md, .py, .js, .ts, など）"""
    
    @staticmethod
    def parse(file_path: str) -> List[Dict[str, str]]:
        """
        テキストファイルを丸ごと読み込む
        
        Args:
            file_path: テキストファイルのパス
        
        Returns:
            ファイル全体のテキストと位置情報のリスト
        """
        content = ""
        encodings = ['utf-8', 'shift-jis', 'cp932', 'euc-jp', 'iso-2022-jp']
        
        # 複数のエンコーディングを試行
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                    break
            except (UnicodeDecodeError, LookupError):
                continue
        
        # すべてのエンコーディングで失敗した場合、エラーハンドリング付きでUTF-8を使用
        if not content:
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
            except Exception:
                # 最後の手段としてバイナリモードで読み込み（chardetがある場合）
                try:
                    import chardet
                    with open(file_path, 'rb') as f:
                        raw_data = f.read()
                        detected = chardet.detect(raw_data)
                        encoding = detected.get('encoding') or 'utf-8'
                        content = raw_data.decode(encoding, errors='ignore')
                except ImportError:
                    # chardetがない場合はエラーを無視
                    pass
        
        if content.strip():
            return [{
                "content": content,
                "location_info": "Full Document"
            }]
        return []


class CodeParser(FileParser):
    """コードファイルパーサー（.py, .ts, .tsx, .js, .jsxなど）"""
    
    # コードファイルの拡張子
    CODE_EXTENSIONS = {'.py', '.ts', '.tsx', '.js', '.jsx', '.java', '.cpp', '.c', '.h', '.hpp', 
                       '.cs', '.go', '.rs', '.rb', '.php', '.swift', '.kt', '.scala', '.r', 
                       '.m', '.mm', '.sh', '.bash', '.zsh', '.fish', '.ps1', '.bat', '.cmd'}
    
    @staticmethod
    def is_code_file(file_path: str) -> bool:
        """ファイルがコードファイルかどうかを判定"""
        ext = Path(file_path).suffix.lower()
        return ext in CodeParser.CODE_EXTENSIONS
    
    @staticmethod
    def parse(file_path: str) -> List[Dict[str, str]]:
        """
        コードファイルを解析してトークンと位置情報を抽出
        
        Args:
            file_path: コードファイルのパス
        
        Returns:
            トークンと位置情報のリスト（通常のパーサーとは異なる形式）
        """
        content = ""
        encodings = ['utf-8', 'shift-jis', 'cp932', 'euc-jp', 'iso-2022-jp']
        
        # 複数のエンコーディングを試行
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                    break
            except (UnicodeDecodeError, LookupError):
                continue
        
        if not content:
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
            except Exception:
                return []
        
        if not content.strip():
            return []
        
        # コードを解析してトークンと位置情報を抽出
        tokens = CodeParser._extract_tokens(content, file_path)
        return tokens
    
    @staticmethod
    def _extract_tokens(content: str, file_path: str) -> List[Dict[str, str]]:
        """
        コードからトークン（識別子、キーワード、文字列リテラルなど）を抽出
        
        Args:
            content: コードの内容
            file_path: ファイルパス
        
        Returns:
            トークン情報のリスト
        """
        tokens = []
        lines = content.split('\n')
        
        # 正規表現パターン
        # 識別子（変数名、関数名、クラス名など）
        identifier_pattern = re.compile(r'\b[a-zA-Z_][a-zA-Z0-9_]*\b')
        # 文字列リテラル
        string_pattern = re.compile(r'["\'](?:[^"\'\\]|\\.)*["\']')
        # 数値リテラル
        number_pattern = re.compile(r'\b\d+\.?\d*\b')
        # コメント（行コメントとブロックコメント）
        comment_pattern = re.compile(r'//.*?$|/\*.*?\*/|#.*?$', re.MULTILINE | re.DOTALL)
        
        for line_num, line in enumerate(lines, start=1):
            # 識別子を抽出
            for match in identifier_pattern.finditer(line):
                token = match.group()
                # キーワードや短すぎるトークンは除外
                if len(token) >= 2 and not CodeParser._is_keyword(token):
                    tokens.append({
                        'token': token.lower(),  # 小文字に正規化
                        'line': line_num,
                        'column': match.start(),
                        'type': 'identifier'
                    })
            
            # 文字列リテラルを抽出
            for match in string_pattern.finditer(line):
                token = match.group()
                # 文字列の中身を抽出（クォートを除く）
                inner = token[1:-1]
                if len(inner) >= 2:
                    tokens.append({
                        'token': inner.lower(),
                        'line': line_num,
                        'column': match.start(),
                        'type': 'string'
                    })
        
        return tokens
    
    @staticmethod
    def _is_keyword(token: str) -> bool:
        """トークンがプログラミング言語のキーワードかどうかを判定"""
        # 主要なプログラミング言語のキーワード（簡易版）
        keywords = {
            'if', 'else', 'elif', 'for', 'while', 'do', 'switch', 'case', 'break', 'continue',
            'return', 'def', 'class', 'import', 'from', 'as', 'try', 'except', 'finally',
            'with', 'async', 'await', 'const', 'let', 'var', 'function', 'export', 'default',
            'public', 'private', 'protected', 'static', 'final', 'abstract', 'interface',
            'extends', 'implements', 'new', 'this', 'super', 'null', 'undefined', 'true', 'false',
            'and', 'or', 'not', 'in', 'is', 'None', 'True', 'False', 'pass', 'yield', 'lambda'
        }
        return token.lower() in keywords


def get_parser(file_path: str) -> Optional[FileParser]:
    """
    ファイルパスに基づいて適切なパーサーを取得
    
    Args:
        file_path: ファイルのパス
    
    Returns:
        対応するパーサークラス、対応していない場合はNone
    """
    ext = Path(file_path).suffix.lower()
    
    # コードファイルの場合はCodeParserを使用
    if CodeParser.is_code_file(file_path):
        return CodeParser()
    
    parsers = {
        '.pdf': PDFParser,
        '.docx': WordParser,
        '.pptx': PowerPointParser,
        '.xlsx': ExcelParser,
        '.txt': TextParser,
        '.md': TextParser,
        '.markdown': TextParser,
        '.json': TextParser,
        '.xml': TextParser,
        '.html': TextParser,
        '.css': TextParser,
        '.yaml': TextParser,
        '.yml': TextParser,
        '.csv': TextParser,
    }
    
    parser_class = parsers.get(ext)
    return parser_class() if parser_class else None


def get_code_parser(file_path: str) -> Optional[CodeParser]:
    """
    コードファイル用のパーサーを取得
    
    Args:
        file_path: コードファイルのパス
    
    Returns:
        CodeParserインスタンス、コードファイルでない場合はNone
    """
    if CodeParser.is_code_file(file_path):
        return CodeParser()
    return None

